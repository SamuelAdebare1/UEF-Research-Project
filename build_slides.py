"""
build_slides.py  —  generates the RAG research presentation as a .pptx file.

Run:
    python build_slides.py
Output:
    RAG_Research_Presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

BASE_DIR = Path(__file__).parent

# ── colour palette ─────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0xFF, 0xFF, 0xFF)   # white background
ACCENT     = RGBColor(0xF0, 0xF4, 0xF8)   # very light blue-gray (card bg)
HIGHLIGHT  = RGBColor(0xDB, 0xEA, 0xF8)   # light blue highlight
TEAL       = RGBColor(0x00, 0x6A, 0xB0)   # UEF blue accent
WHITE      = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy — used for text (inverted)
LIGHT_GRAY = RGBColor(0x44, 0x44, 0x55)   # dark gray for body text
YELLOW     = RGBColor(0xC0, 0x6A, 0x00)   # dark amber (readable on white)
GREEN      = RGBColor(0x18, 0x7A, 0x3B)   # dark green
RED        = RGBColor(0xC0, 0x2A, 0x1A)   # dark red
ORANGE     = RGBColor(0xC0, 0x5A, 0x00)   # dark orange

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]   # truly blank


# ── helper functions ───────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # full-bleed dark background
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_header_bar(slide, title, subtitle=None):
    """Adds the teal left accent bar + title text."""
    # left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, TEAL)
    # title
    add_text(slide, title,
             Inches(0.3), Inches(0.18), Inches(12.7), Inches(0.7),
             font_size=Pt(32), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.3), Inches(0.82), Inches(12.7), Inches(0.4),
                 font_size=Pt(16), color=TEAL)
    # thin divider line
    add_rect(slide, Inches(0.3), Inches(1.18), Inches(12.73), Pt(2), TEAL)


def add_bullet_box(slide, items, x, y, w, h,
                   title=None, title_color=TEAL,
                   bullet="•", font_size=Pt(17), line_gap=Inches(0.38)):
    """Adds a multi-line bullet list as a textbox."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.bold = True
        r.font.size = Pt(19)
        r.font.color.rgb = title_color
        p = tf.add_paragraph()
        p.text = ""

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = font_size
        r.font.color.rgb = WHITE
        p.space_after = Pt(6)


def card(slide, x, y, w, h, fill=ACCENT, border_color=TEAL, border_w=Pt(1.5)):
    return add_rect(slide, x, y, w, h, fill, border_color, border_w)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

# background accent
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, TEAL)
add_rect(s, Inches(0), Inches(5.5), SLIDE_W, Inches(2.0), HIGHLIGHT)

add_text(s, "Information Loss in RAG Systems",
         Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.2),
         font_size=Pt(44), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s, "A Needle-in-a-Haystack Evaluation of Retrieval-Augmented Generation",
         Inches(0.5), Inches(2.6), Inches(11), Inches(0.7),
         font_size=Pt(22), color=TEAL, align=PP_ALIGN.LEFT)

add_text(s, "Samuel Oyewoade  ·  May 2026",
         Inches(0.5), Inches(5.7), Inches(12), Inches(0.5),
         font_size=Pt(16), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

add_text(s, "",
         Inches(0.5), Inches(6.15), Inches(12), Inches(0.4),
         font_size=Pt(15), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — RESEARCH QUESTION & MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Research Question", "Why does RAG lose information — and how much?")

# big question box
c = card(s, Inches(0.35), Inches(1.35), Inches(12.63), Inches(1.4))
add_text(s, "When a language model answers questions from a long document using RAG\n"
            "instead of the full context, how much information is lost — and where does it go?",
         Inches(0.55), Inches(1.45), Inches(12.2), Inches(1.2),
         font_size=Pt(20), color=YELLOW, align=PP_ALIGN.LEFT, bold=True)

add_bullet_box(s, [
    "Modern LLMs have context limits — a 50-page document can exhaust a smaller model's window",
    "RAG is the dominant industry solution: retrieve only what's relevant",
    "But retrieval introduces a new failure mode — the answer may exist in the document but not in the retrieved chunk",
    "We quantify this gap by comparing full-context accuracy vs. RAG accuracy on identical questions",
], Inches(0.35), Inches(2.9), Inches(12.6), Inches(3.5), font_size=Pt(18))

add_text(s, "Two conditions  →  one measure of information loss",
         Inches(0.35), Inches(6.65), Inches(12.6), Inches(0.5),
         font_size=Pt(17), color=TEAL, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE DATA
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "The Data", "50-page document + controlled needle injections")

# left column: document description
card(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(4.8))
add_text(s, "Source Document",
         Inches(0.55), Inches(1.45), Inches(5.7), Inches(0.45),
         font_size=Pt(20), bold=True, color=TEAL)
add_bullet_box(s, [
    "Base text: Book of Genesis (KJV) — highly consistent, archaic prose",
    "50 pages total  ·  ~35,000 tokens  ·  113 chunks after splitting",
    "5 fictional 'needle' passages injected at pages 1, 15, 25, 35, 50",
    "Needles span 5 different fictional genres (compliance manual, research paper, financial contract, space telemetry, historical archive)",
    "Purpose: facts unknown to any LLM from pre-training — retrieval is the ONLY way to answer correctly",
], Inches(0.55), Inches(1.95), Inches(5.7), Inches(3.8), font_size=Pt(15))

# right column: question set
card(s, Inches(6.65), Inches(1.35), Inches(6.33), Inches(4.8))
add_text(s, "Evaluation Set — 30 Questions",
         Inches(6.85), Inches(1.45), Inches(6.0), Inches(0.45),
         font_size=Pt(20), bold=True, color=TEAL)

rows = [
    ("Q1–Q5",   "Needle",           "Injected facts only. Correct answer requires retrieval — cannot be guessed."),
    ("Q6–Q15",  "Adjacent-to-Needle","Details from the same chunk as a needle. Tests context richness."),
    ("Q16–Q30", "Biblical Fact",     "Genesis facts the LLM may know from pre-training — tests retrieval independence."),
]
y = Inches(1.95)
for qrange, qtype, qdesc in rows:
    add_text(s, f"{qrange}  —  {qtype}",
             Inches(6.85), y, Inches(6.0), Inches(0.35),
             font_size=Pt(15), bold=True, color=ORANGE)
    add_text(s, qdesc,
             Inches(6.85), y + Inches(0.33), Inches(6.0), Inches(0.5),
             font_size=Pt(13), color=LIGHT_GRAY)
    y += Inches(0.95)

add_text(s, "Difficulty tags: easy / medium / hard  ·  Section tags: early / middle / late",
         Inches(6.85), Inches(5.5), Inches(6.0), Inches(0.4),
         font_size=Pt(13), color=LIGHT_GRAY, italic=True)

# needle table at bottom
card(s, Inches(0.35), Inches(6.2), Inches(12.63), Inches(0.9))
add_text(s,
         "Needles at depth:  Page 1 → Wristband  ·  Page 15 → Metallic taste  ·  Page 25 → Helsinki  ·  Page 35 → Xylanthia-9  ·  Page 50 → Secondary rudder",
         Inches(0.55), Inches(6.28), Inches(12.3), Inches(0.7),
         font_size=Pt(14), color=YELLOW, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHY RAG + COMPLEXITIES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Why RAG? The Complexities", "Long documents, limited context, and retrieval trade-offs")

# three-column complexity cards
cols = [
    ("Problem: Context Limits",
     HIGHLIGHT,
     [
         "A 50-page document (~23k tokens) exceeds many model context windows",
         "Even 128k-context models degrade on long-range recall — the 'lost-in-the-middle' effect",
         "Feeding the full document is slow, expensive, and does not scale to enterprise corpora",
     ]),
    ("Solution: RAG",
     RGBColor(0xD4, 0xF0, 0xE2),
     [
         "Retrieve only the top-k most relevant chunks at query time",
         "LLM sees a focused ~5-chunk context → faster, cheaper, scalable",
         "Dominant pattern in production NLP systems (chatbots, document Q&A, legal search)",
     ]),
    ("New Problem: Information Loss",
     RGBColor(0xFA, 0xE2, 0xE2),
     [
         "Chunking splits continuous text → a fact may span two chunks, neither complete",
         "Retrieval is imperfect — the right chunk may rank 6th when top-5 are returned",
         "Even when the right chunk IS retrieved, the LLM may ignore it and use memorised knowledge",
     ]),
]

x = Inches(0.35)
for col_title, col_color, bullets in cols:
    c = card(s, x, Inches(1.35), Inches(4.1), Inches(5.4), fill=col_color, border_color=TEAL)
    add_text(s, col_title,
             x + Inches(0.15), Inches(1.45), Inches(3.8), Inches(0.5),
             font_size=Pt(16), bold=True, color=YELLOW)
    add_bullet_box(s, bullets,
                   x + Inches(0.15), Inches(1.95), Inches(3.8), Inches(4.5),
                   font_size=Pt(14))
    x += Inches(4.31)

add_text(s, "Our experiment isolates the retrieval layer: same document, same questions, same models — only the context delivery changes.",
         Inches(0.35), Inches(6.85), Inches(12.63), Inches(0.5),
         font_size=Pt(15), color=TEAL, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RAG ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "RAG Architecture", "Three-stage pipeline: chunk → embed → retrieve → generate")

# pipeline boxes
stages = [
    ("1  PDF Ingestion", "50-pages.pdf\n(PyMuPDF)\nRaw text\nextracted"),
    ("2  Chunking", "chunker.py\nSliding window\n350 tok / 35 overlap\n→ chunks.json\n(113 chunks)"),
    ("3  Embedding", "embedder.py\nall-mpnet-base-v2\n768-dim vectors\n→ embeddings.json"),
    ("4  Query Time", "User question\nEmbedded with\nsame model\n↓\nCosine similarity"),
    ("5  Generation", "Top-5 chunks\n+ question\n→ Llama / Mistral\n→ Answer"),
]
colors = [HIGHLIGHT, RGBColor(0xD4, 0xF0, 0xE2), RGBColor(0xF0, 0xE4, 0xFA),
          RGBColor(0xFA, 0xF0, 0xD4), RGBColor(0xFA, 0xE2, 0xE2)]

box_w = Inches(2.35)
box_h = Inches(3.6)
gap   = Inches(0.18)
start_x = Inches(0.28)
y_box = Inches(1.55)

for i, ((title, body), col) in enumerate(zip(stages, colors)):
    bx = start_x + i * (box_w + gap)
    card(s, bx, y_box, box_w, box_h, fill=col, border_color=TEAL, border_w=Pt(1.5))
    add_text(s, title, bx + Inches(0.1), y_box + Inches(0.1),
             box_w - Inches(0.2), Inches(0.45),
             font_size=Pt(15), bold=True, color=YELLOW)
    add_text(s, body, bx + Inches(0.1), y_box + Inches(0.55),
             box_w - Inches(0.2), Inches(2.9),
             font_size=Pt(13), color=WHITE)
    # arrow (except after last box)
    if i < 4:
        ax = bx + box_w + Inches(0.01)
        add_text(s, "→", ax, y_box + Inches(1.5),
                 Inches(0.18), Inches(0.45),
                 font_size=Pt(22), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# local LLM note
card(s, Inches(0.28), Inches(5.3), Inches(12.77), Inches(0.75), fill=ACCENT, border_color=TEAL)
add_text(s,
         "Fully local pipeline  ·  No API key required  ·  LLMs served via GPT4All (Llama 3.1 8B 128k  &  Mistral 7B Instruct)  ·  Embeddings via sentence-transformers",
         Inches(0.5), Inches(5.38), Inches(12.4), Inches(0.6),
         font_size=Pt(14), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s,
         "API layer: FastAPI  ·  /query (batch)  ·  /query/stream (SSE token-by-token)  ·  React chat UI",
         Inches(0.5), Inches(6.1), Inches(12.4), Inches(0.5),
         font_size=Pt(14), color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CHUNKING STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Chunking Strategy", "How and why the document is split")

# left: parameters table-style
card(s, Inches(0.35), Inches(1.35), Inches(5.8), Inches(5.5))
add_text(s, "Parameters", Inches(0.55), Inches(1.45),
         Inches(5.5), Inches(0.4), font_size=Pt(20), bold=True, color=TEAL)

params = [
    ("Chunk size",   "350 tokens",     "Focused context per chunk; reduces noise fed to LLM"),
    ("Overlap",      "35 tokens (10%)", "Prevents a needle from splitting across two chunk boundaries"),
    ("Tokenizer",    "cl100k_base",    "Same as GPT-4 / OpenAI embedding models — reproducible sizing"),
    ("Chunking tool","tiktoken",       "Window sizing only — embedding uses its own internal tokenizer"),
    ("Output",       "113 chunks",     "Avg 350 tokens each (~1,300 chars)"),
]
y = Inches(1.95)
for pname, pval, preason in params:
    add_text(s, pname, Inches(0.55), y, Inches(1.5), Inches(0.35),
             font_size=Pt(13), color=ORANGE, bold=True)
    add_text(s, pval, Inches(2.1), y, Inches(1.6), Inches(0.35),
             font_size=Pt(13), color=WHITE, bold=True)
    add_text(s, preason, Inches(0.55), y + Inches(0.33), Inches(5.3), Inches(0.35),
             font_size=Pt(12), color=LIGHT_GRAY)
    y += Inches(0.85)

# right: sliding window diagram (text art)
card(s, Inches(6.45), Inches(1.35), Inches(6.53), Inches(5.5))
add_text(s, "Sliding Window Logic", Inches(6.65), Inches(1.45),
         Inches(6.2), Inches(0.4), font_size=Pt(20), bold=True, color=TEAL)

diagram = (
    "Token stream (full document):\n"
    "│ t₀  t₁  t₂  … t₃₄₉ │ t₃₅₀ … │\n\n"
    "Chunk 0:  ├──────────────────────┤\n"
    "           [tokens 0 … 349]\n\n"
    "Chunk 1:             ├──────────────────────┤\n"
    "                     [tokens 315 … 664]\n"
    "                     └──35-token overlap──┘\n\n"
    "Chunk 2:                          ├────────────…\n"
    "                                  [tokens 630 … ]\n\n"
    "step = chunk_size − overlap = 315 tokens\n"
    "→ 113 chunks for the 50-page document"
)
add_text(s, diagram, Inches(6.65), Inches(1.95), Inches(6.2), Inches(4.5),
         font_size=Pt(13), color=WHITE, italic=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EMBEDDING & RETRIEVAL
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Embedding & Retrieval", "Turning text into searchable vectors")

# embed model card
card(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(2.6))
add_text(s, "Embedding Model: all-mpnet-base-v2",
         Inches(0.55), Inches(1.45), Inches(5.8), Inches(0.45),
         font_size=Pt(18), bold=True, color=TEAL)
props = [
    "Vector dimension : 768",
    "Max sequence length : 384 tokens",
    "Training data : 1 billion+ sentence pairs",
    "Why chosen : best general-quality model in SBERT catalogue",
    "Runs fully offline after first download",
]
add_bullet_box(s, props, Inches(0.55), Inches(1.95), Inches(5.7), Inches(2.3),
               font_size=Pt(14))

# cosine similarity card
card(s, Inches(0.35), Inches(4.1), Inches(6.1), Inches(2.0))
add_text(s, "Retrieval: Cosine Similarity",
         Inches(0.55), Inches(4.2), Inches(5.8), Inches(0.45),
         font_size=Pt(18), bold=True, color=TEAL)
add_bullet_box(s, [
    "Query embedded with the same model at run-time",
    "Cosine similarity against all 113 chunk vectors",
    "Top-k = 5 chunks returned to the LLM",
    "No approximate search needed (79 vectors, exact is fast)",
], Inches(0.55), Inches(4.7), Inches(5.7), Inches(1.8), font_size=Pt(14))

# formula
card(s, Inches(6.65), Inches(1.35), Inches(6.33), Inches(2.6))
add_text(s, "Cosine Similarity Formula",
         Inches(6.85), Inches(1.45), Inches(6.0), Inches(0.4),
         font_size=Pt(18), bold=True, color=TEAL)
add_text(s,
         "sim(q, cᵢ) = (q · cᵢ) / (‖q‖ · ‖cᵢ‖)\n\n"
         "q   = query embedding vector\n"
         "cᵢ  = chunk i embedding vector\n\n"
         "Ranks by semantic similarity, not keyword overlap\n"
         "→ handles paraphrase and synonymy naturally",
         Inches(6.85), Inches(1.9), Inches(6.0), Inches(2.0),
         font_size=Pt(14), color=WHITE)

# LLM config
card(s, Inches(6.65), Inches(4.1), Inches(6.33), Inches(2.0))
add_text(s, "LLM Configuration",
         Inches(6.85), Inches(4.2), Inches(6.0), Inches(0.4),
         font_size=Pt(18), bold=True, color=TEAL)
add_bullet_box(s, [
    "Llama 3.1 8B Instruct (128k context, Q4_0 quantized)",
    "Mistral 7B Instruct (32k context, Q4_0 quantized)",
    "Context window: 4096 tokens per query (hardware limit)",
    "Temperature: 0.1 (deterministic recall answers)",
], Inches(6.85), Inches(4.7), Inches(6.0), Inches(1.8), font_size=Pt(14))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — EXPERIMENT DESIGN
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Experiment Design", "Needle-in-a-Haystack controlled evaluation")

# two conditions side by side
card(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(3.3), fill=RGBColor(0xD4, 0xF0, 0xE2))
add_text(s, "Condition A — Full Context (Baseline)",
         Inches(0.55), Inches(1.45), Inches(5.7), Inches(0.4),
         font_size=Pt(17), bold=True, color=GREEN)
add_bullet_box(s, [
    "Different document: story 52845 from the QuALITY dataset (fiction, not 50-pages.pdf)",
    "22 questions  ·  3 needle injections (not 5)",
    "Each question asked manually in a fresh, isolated session",
    "Models: DeepSeek-R1-14B, Llama 3.1 8B, Mistral 7B",
    "Accuracy: Llama 61.4%  ·  Mistral 45.5%  ·  DeepSeek 40.9%",
    "Note: different document — not a direct apples-to-apples comparison with Condition B",
], Inches(0.55), Inches(1.9), Inches(5.7), Inches(2.5), font_size=Pt(13))

card(s, Inches(6.65), Inches(1.35), Inches(6.33), Inches(3.3), fill=HIGHLIGHT)
add_text(s, "Condition B — RAG (Our System)",
         Inches(6.85), Inches(1.45), Inches(6.0), Inches(0.4),
         font_size=Pt(17), bold=True, color=TEAL)
add_bullet_box(s, [
    "Only top-5 retrieved chunks shown to the LLM",
    "Automated test runner — 30 questions × 2 models",
    "Models: Llama 3.1 8B 128k, Mistral 7B Instruct",
    "Accuracy: Llama 93.3%  ·  Mistral 86.7%",
], Inches(6.85), Inches(1.9), Inches(6.0), Inches(2.5), font_size=Pt(14))

# notable result highlight
card(s, Inches(0.35), Inches(4.75), Inches(12.63), Inches(0.85), fill=RGBColor(0xFF, 0xF0, 0xCC))
add_text(s,
         "Note: Condition A used a different document (QuALITY story 52845). "
         "Direct comparison is indicative, not controlled — RAG accuracy is higher, but documents differ.",
         Inches(0.55), Inches(4.83), Inches(12.3), Inches(0.7),
         font_size=Pt(15), color=YELLOW, bold=True)

# needle design
card(s, Inches(0.35), Inches(5.7), Inches(12.63), Inches(1.5))
add_text(s, "Needle Placement Design",
         Inches(0.55), Inches(5.78), Inches(3.5), Inches(0.35),
         font_size=Pt(16), bold=True, color=TEAL)
needles = [
    "Page 1\nWristband",
    "Page 15\nMetallic taste",
    "Page 25\nHelsinki",
    "Page 35\nXylanthia-9",
    "Page 50\nSecondary rudder",
]
nx = Inches(0.5)
for n in needles:
    card(s, nx, Inches(6.15), Inches(2.35), Inches(0.8),
         fill=RGBColor(0xDB, 0xEA, 0xF8), border_color=ORANGE)
    add_text(s, n, nx + Inches(0.1), Inches(6.2), Inches(2.15), Inches(0.7),
             font_size=Pt(11), color=ORANGE, align=PP_ALIGN.CENTER)
    nx += Inches(2.55)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Results", "RAG accuracy by model, category, and document section")

# metrics table
headers = ["Metric", "Llama 3.1 8B 128k", "Mistral 7B Instruct"]
rows_data = [
    ("Overall Accuracy (30 questions)", "93.3%  (28/30)", "86.7%  (26/30)"),
    ("Needle Score  (Q1–Q5)",           "100.0%  (5/5)",  "100.0%  (5/5)"),
    ("Adjacent-to-Needle  (Q6–Q15)",    "90.0%  (9/10)",  "100.0%  (10/10)"),
    ("Fact Recall  (Q16–Q30)",          "93.3%  (14/15)", "73.3%  (11/15)"),
    ("Section — early  (pp 1–15)",      "93.3%  (14/15)", "80.0%  (12/15)"),
    ("Section — middle (pp 16–33)",     "100.0%  (4/4)",  "100.0%  (4/4)"),
    ("Section — late  (pp 34–50)",      "90.9%  (10/11)", "90.9%  (10/11)"),
    ("Difficulty — easy",               "100.0%  (4/4)",  "100.0%  (4/4)"),
    ("Difficulty — medium",             "94.4%  (17/18)", "83.3%  (15/18)"),
    ("Difficulty — hard",               "87.5%  (7/8)",   "87.5%  (7/8)"),
]

# table header
add_rect(s, Inches(0.35), Inches(1.35), Inches(12.63), Inches(0.5), TEAL)
col_xs = [Inches(0.45), Inches(5.6), Inches(9.35)]
col_ws = [Inches(5.0), Inches(3.6), Inches(3.5)]
for i, hdr in enumerate(headers):
    add_text(s, hdr, col_xs[i], Inches(1.38), col_ws[i], Inches(0.45),
             font_size=Pt(14), bold=True, color=DARK_BG, align=PP_ALIGN.LEFT)

row_bg_a = RGBColor(0xF0, 0xF0, 0xF8)
row_bg_b = ACCENT
y = Inches(1.85)
for ri, (metric, llama_v, mistral_v) in enumerate(rows_data):
    bg = row_bg_a if ri % 2 == 0 else row_bg_b
    add_rect(s, Inches(0.35), y, Inches(12.63), Inches(0.45), bg)
    add_text(s, metric, col_xs[0], y + Inches(0.04), col_ws[0], Inches(0.4),
             font_size=Pt(13), color=WHITE)
    # colour code the values
    for val, cx, cw in [(llama_v, col_xs[1], col_ws[1]), (mistral_v, col_xs[2], col_ws[2])]:
        pct_val = float(val.split('%')[0])
        vc = GREEN if pct_val >= 90 else (ORANGE if pct_val >= 80 else RED)
        add_text(s, val, cx, y + Inches(0.04), cw, Inches(0.4),
                 font_size=Pt(13), bold=True, color=vc)
    y += Inches(0.45)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — UI SCREENSHOT: EMPTY STATE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "RAG Chat UI — Landing Screen",
               "React frontend connected to the local FastAPI + Llama 3.1 8B backend")
# UI screenshot: 1400×860, fits well at full width below header
ui_w = Inches(12.0)
ui_h = Inches(12.0) * (860 / 1400)
ui_x = (SLIDE_W - ui_w) / 2
s.shapes.add_picture(str(BASE_DIR / "ui_empty.png"), ui_x, Inches(1.3), ui_w, ui_h)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — UI SCREENSHOT: QUERY IN ACTION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "RAG Chat UI — Query & Retrieved Sources",
               "Top-5 chunks ranked by cosine similarity, answer streaming token-by-token")
s.shapes.add_picture(str(BASE_DIR / "ui_answer_done.png"), ui_x, Inches(1.3), ui_w, ui_h)
# annotation callout
add_text(s, "Chunk 54 = 36.0% similarity → correct needle chunk  ·  Answer: Xylanthia-9  ✓",
         Inches(0.35), Inches(7.0), Inches(12.63), Inches(0.38),
         font_size=Pt(14), color=RGBColor(0x00, 0x6A, 0xB0), bold=True,
         align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — FULL-CONTEXT BASELINE vs RAG
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Full-Context Baseline vs RAG", "Same models, same questions — only context delivery changes")

# comparison table header
add_rect(s, Inches(0.35), Inches(1.35), Inches(12.63), Inches(0.5), TEAL)
for txt, cx in [("Metric", Inches(0.5)), ("DeepSeek-R1-14B", Inches(3.8)),
                ("Llama 3.1 8B", Inches(6.6)), ("Mistral 7B", Inches(9.4))]:
    add_text(s, txt, cx, Inches(1.38), Inches(2.7), Inches(0.45),
             font_size=Pt(14), bold=True, color=DARK_BG)

full_ctx_rows = [
    ("Overall Accuracy",       "40.9%", "61.4%", "45.5%"),
    ("Needle Retrieval",       "100.0%","100.0%","66.7%"),
    ("Section — early",        "37.5%", "50.0%", "50.0%"),
    ("Section — middle",       "33.3%", "61.1%", "33.3%"),
    ("Section — late",         "60.0%", "80.0%", "60.0%"),
]
rag_rows = [
    ("—",      "93.3%", "86.7%"),
    ("—",      "100.0%","100.0%"),
    ("—",      "93.3%", "80.0%"),
    ("—",      "100.0%","100.0%"),
    ("—",      "90.9%", "90.9%"),
]

y = Inches(1.85)
for ri, (row_fc, row_rag) in enumerate(zip(full_ctx_rows, rag_rows)):
    bg = RGBColor(0xF0, 0xF0, 0xF8) if ri % 2 == 0 else ACCENT
    add_rect(s, Inches(0.35), y, Inches(12.63), Inches(0.5), bg)
    add_text(s, row_fc[0], Inches(0.5), y + Inches(0.06), Inches(3.1), Inches(0.4),
             font_size=Pt(13), color=WHITE)
    # full-context values
    for val, cx in zip(row_fc[1:], [Inches(3.8), Inches(6.6), Inches(9.4)]):
        pv = float(val.replace('%',''))
        vc = GREEN if pv >= 80 else (ORANGE if pv >= 60 else RED)
        add_text(s, f"FC: {val}", cx, y + Inches(0.06), Inches(2.5), Inches(0.4),
                 font_size=Pt(12), color=vc)
    y += Inches(0.5)

# RAG overlay note
card(s, Inches(0.35), Inches(4.6), Inches(12.63), Inches(1.05), fill=RGBColor(0xD4, 0xF0, 0xE2))
add_text(s, "RAG Condition Results (same models):",
         Inches(0.55), Inches(4.67), Inches(4.0), Inches(0.35),
         font_size=Pt(14), bold=True, color=GREEN)
add_text(s, "DeepSeek: failed to load via GPT4All    Llama: 93.3%  (+31.9 pp)    Mistral: 86.7%  (+41.2 pp)",
         Inches(0.55), Inches(4.99), Inches(12.2), Inches(0.45),
         font_size=Pt(15), bold=True, color=YELLOW)

# big insight box
card(s, Inches(0.35), Inches(5.75), Inches(12.63), Inches(1.35), fill=RGBColor(0xE4, 0xF5, 0xD4))
add_text(s, "Why RAG outperforms full-context here:",
         Inches(0.55), Inches(5.82), Inches(5.5), Inches(0.4),
         font_size=Pt(15), bold=True, color=GREEN)
add_bullet_box(s, [
    "Full context = 35,000 tokens → models lose focus mid-document ('lost in the middle' effect)",
    "RAG = 5 focused chunks (~2,000 tokens) → model sees exactly what it needs, nothing else",
    "Trade-off: RAG can fail if the right chunk is not retrieved — but at top-k=5 it never did for needle questions",
], Inches(0.55), Inches(6.22), Inches(12.2), Inches(0.8), font_size=Pt(13))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — APP OUTPUT: SUCCESSFUL RAG ANSWERS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "App Output — RAG in Action", "Real answers from the test runner (30 questions × 2 models)")

examples = [
    ("Q1 — Needle (page 1)",
     "Which wristband must employees wear when operating TS-500 machinery in Sector 4?",
     "Mistral: Employees operating the TS-500 machinery in Sector 4 must wear\na neon-orange anti-static wristband on their left arm at all times.",
     "Neon-orange anti-static wristband", GREEN, "Retrieved chunks: [1, 7, 52, 70, 75]  ✓ Chunk 1 contains the needle"),
    ("Q3 — Needle (page 25)",
     "To which city must the Request for Dissolution form be submitted?",
     'Mistral: The Request for Dissolution form should be submitted directly\nto the managing director\'s office in Helsinki.',
     "Helsinki", GREEN, "Retrieved chunks: [42, 43, 44, 45, 46]  ✓ Chunk 43 contains the needle"),
    ("Q4 — Needle (page 35)",
     "What is the name of the asteroid near which the secondary probe detected a temperature drop?",
     "Llama: Xylanthia-9.",
     "Xylanthia-9", GREEN, "Retrieved chunks: [54, 55, 56, 57, 58]  ✓ Chunk 54 contains the needle"),
    ("Q30 — Biblical Fact",
     "What did Rebekah instruct Jacob to fetch from the flock?",
     "Mistral: Two good kids of the goats.",
     "Two good kids of the goats", GREEN, "Retrieved chunk confirmed — correct and grounded in context"),
]

y = Inches(1.38)
for qid, question, answer, target, col, note in examples:
    card(s, Inches(0.35), y, Inches(12.63), Inches(1.32), fill=RGBColor(0xD4, 0xF0, 0xE2), border_color=col)
    add_text(s, qid, Inches(0.52), y + Inches(0.07), Inches(2.5), Inches(0.32),
             font_size=Pt(13), bold=True, color=col)
    add_text(s, f"Q: {question}", Inches(0.52), y + Inches(0.35), Inches(12.1), Inches(0.32),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)
    add_text(s, answer, Inches(0.52), y + Inches(0.65), Inches(9.5), Inches(0.42),
             font_size=Pt(13), color=WHITE, bold=True)
    add_text(s, note, Inches(0.52), y + Inches(1.05), Inches(12.1), Inches(0.25),
             font_size=Pt(11), color=TEAL)
    y += Inches(1.42)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — APP OUTPUT: FAILURE MODE EXAMPLES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "App Output — Failure Modes", "Real wrong answers and what they reveal")

failures = [
    ("Q11 — Llama  [Within-chunk attention failure]",
     "Which arm of the galaxy was the probe scanning?",
     'Llama: "I don\'t know based on the document."',
     "Perseus arm",
     "Chunk 54 WAS retrieved (same chunk that answered Q4=Xylanthia-9 correctly). Llama found Xylanthia-9\n3 sentences later but missed 'Perseus arm' earlier in the same chunk — attention collapsed on a sub-passage.",
     ORANGE),
    ("Q21 — Mistral  [Partial-context read]",
     "How many total years did Methuselah live?",
     '"Methuselah lived for a total of 187 years."',
     "969 years",
     "187 is Methuselah's age when he BEGAT Lamech — a number two sentences before the total. Mistral\nextracted the first relevant number it encountered rather than reading to the total-lifespan sentence.",
     ORANGE),
    ("Q23 — Mistral  [LLM assumption overrides retrieved context]",
     "What did God promise Abram at the plain of Moreh in Sichem?",
     '"God promised Abram the land from the river of Egypt to the Euphrates."',
     '"Unto thy seed will I give this land" (Genesis 12:7)',
     "Retrieved chunk contained Genesis 12:7 (the Sichem promise). Mistral answered with Genesis 15:18 —\na different, larger covenant NOT in the retrieved context. Pre-training knowledge overrode the retrieval.",
     RED),
    ("Q27 — Mistral  [Training memory beats retrieved passage]",
     "How old was Abraham when he died?",
     '"Abraham died at an age of one hundred and thirty-seven years."',
     "175 years (Genesis 25:7)",
     "137 is Sarah's age at death (Genesis 23:1). Chunk 54 containing Abraham's 175-year figure WAS retrieved.\nMistral produced a plausible biblical number from memorised training data rather than reading the chunk.",
     RED),
]

y = Inches(1.38)
for qid, question, model_ans, correct, analysis, col in failures:
    card(s, Inches(0.35), y, Inches(12.63), Inches(1.42), fill=RGBColor(0xFA, 0xE2, 0xE2), border_color=col)
    add_text(s, qid, Inches(0.52), y + Inches(0.07), Inches(12.1), Inches(0.3),
             font_size=Pt(13), bold=True, color=col)
    add_text(s, f"Q: {question}", Inches(0.52), y + Inches(0.35), Inches(9.0), Inches(0.28),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)
    add_text(s, model_ans, Inches(0.52), y + Inches(0.62), Inches(9.0), Inches(0.28),
             font_size=Pt(12), color=RED, bold=True)
    add_text(s, f"Correct: {correct}", Inches(9.8), y + Inches(0.35), Inches(3.0), Inches(0.28),
             font_size=Pt(12), color=GREEN, bold=True)
    add_text(s, analysis, Inches(0.52), y + Inches(0.92), Inches(12.1), Inches(0.42),
             font_size=Pt(11), color=LIGHT_GRAY)
    y += Inches(1.52)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — FACTS, CONCEPTS & RELATIONS (LINGUISTIC ANALYSIS)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Extracting Facts, Concepts & Relations", "Linguistic analysis of the source document")

# two analysis boxes
analyses = [
    ("Semantic Anomaly Detection",
     "Model: all-MiniLM-L6-v2\n"
     "Threshold: cosine similarity < 0.3\n\n"
     "All 5 injected needles flagged:\n"
     "  • Wristband       (sim = -0.087)\n"
     "  • Helsinki        (sim = -0.053)\n"
     "  • Metallic taste  (sim = +0.052)\n"
     "  • Xylanthia-9     (sim = +0.062)\n"
     "  • Windsor Merchant (sim = +0.081)\n\n"
     "Every injected sentence sits far outside\n"
     "the similarity range of genuine biblical text.\n\n"
     "→ Style-shift detection works;\n"
     "  confirms all 5 injections are detectable",
     ORANGE),
    ("Word Frequency & Vocabulary",
     "Total tokens  : 23,863\n"
     "Unique types  : 1,922\n"
     "TTR (richness): 0.0805\n"
     "Hapax legomena: 839 words (43.7% of vocab)\n\n"
     "Top concepts by frequency:\n"
     "  Abraham (126)  ·  son (122)\n"
     "  earth (113)    ·  Jacob (101)\n"
     "  father (92)    ·  years (84)\n"
     "  Isaac (70)     ·  begat (67)\n\n"
     "→ Frequency map identifies the most\n"
     "  important entities for relation extraction",
     GREEN),
]

x = Inches(0.35)
for title, body, accent_col in analyses:
    card(s, x, Inches(1.35), Inches(6.2), Inches(5.5))
    add_text(s, title, x + Inches(0.15), Inches(1.45),
             Inches(5.9), Inches(0.45),
             font_size=Pt(18), bold=True, color=accent_col)
    add_text(s, body, x + Inches(0.15), Inches(1.95),
             Inches(5.9), Inches(4.6),
             font_size=Pt(14), color=WHITE)
    x += Inches(6.43)

add_text(s, "Relations (next step): co-occurrence + dependency parsing → triplets (Abraham, father-of, Isaac) for knowledge graph construction",
         Inches(0.35), Inches(6.95), Inches(12.63), Inches(0.4),
         font_size=Pt(13), color=TEAL, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — LINGUISTIC ANALYSIS DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Linguistic Analysis Dashboard",
               "Interactive dashboard — anomaly detection & vocabulary profiling")
dash_w = Inches(12.4)
dash_h = dash_w * (1547 / 1400)   # actual screenshot aspect ratio
if dash_h > Inches(5.9):
    dash_h = Inches(5.9)
    dash_w = dash_h * (1400 / 1547)
dash_x = (SLIDE_W - dash_w) / 2
s.shapes.add_picture(str(BASE_DIR / "dashboard_screenshot.png"),
                     dash_x, Inches(1.3), dash_w, dash_h)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — KEY FINDINGS & FAILURE MODES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Key Findings & Failure Modes", "Where information loss actually happens")

findings = [
    (GREEN,  "Finding 1 — Retrieval is NOT the bottleneck",
     "Both models retrieved the correct chunk for all 5 needles (100% needle retrieval). "
     "The all-mpnet-base-v2 embedder with cosine similarity and top-k=5 is reliable at this corpus size. "
     "Information loss does NOT happen at the retrieval step."),
    (ORANGE, "Finding 2 — Within-chunk attention fails on dense passages",
     "Llama retrieved chunk 54 (which contained the Perseus arm answer) yet answered 'I don't know'. "
     "Mistral had the fig-leaf sentence 3 words before 'aprons' and answered 'aprons'. "
     "The model sees the right chunk but fails to extract the precise sentence — a within-context attention gap."),
    (RED,    "Finding 3 — Pre-training knowledge overrides retrieved context",
     "Mistral's Q23: the retrieved chunk contained the Sichem promise (Genesis 12:7). "
     "Mistral answered with the Genesis 15:18 covenant — a passage NOT in the retrieved context. "
     "The model overrode the retrieved evidence with memorised training knowledge. "
     "This is the core failure mode — LLMs assume rather than retrieve."),
    (TEAL,   "Finding 4 — Document depth has no effect on retrieval",
     "Accuracy was consistent across early, middle, and late sections. "
     "Dense embeddings are equally effective at any depth in a 50-page document (113 chunks). "
     "The 'lost-in-the-middle' effect (Attention) does not affect the dense retrieval step."),
]

y = Inches(1.35)
for col, title, body in findings:
    add_rect(s, Inches(0.35), y, Inches(0.1), Inches(0.95), col)
    add_text(s, title, Inches(0.6), y + Inches(0.03),
             Inches(11.8), Inches(0.38),
             font_size=Pt(15), bold=True, color=col)
    add_text(s, body, Inches(0.6), y + Inches(0.38),
             Inches(12.1), Inches(0.55),
             font_size=Pt(13), color=LIGHT_GRAY)
    y += Inches(1.1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header_bar(s, "Next Steps", "Cloud deployment + remaining analyses")

next_steps = [
    ("Cloud Deployment  (to be decided today)",
     [
         "Migrate the FastAPI + GPT4All pipeline to a cloud environment",
         "Options: Google Colab, HuggingFace Spaces, AWS/Azure",
         "Goal: shareable demo, no local hardware dependency",
         "Enable larger models (34B+) that exceed local Apple Silicon memory",
     ], TEAL),
    ("Linguistic Analysis — Remaining",
     [
         "Stylometric analysis: detect injection points via sentence-length & vocabulary-richness shifts",
         "Readability scoring: Flesch-Kincaid on Genesis vs. injected fictional passages",
         "Relation extraction: dependency parsing → (subject, relation, object) triplets",
         "Knowledge graph: entity nodes + relation edges → visual graph of Genesis narrative",
     ], ORANGE),
    ("RAG Improvements",
     [
         "Re-rank retrieved chunks (cross-encoder) to improve within-context extraction",
         "Experiment with chunk sizes: 200 vs 350 vs 500 tokens",
         "Add metadata filters (page number, section) to the retrieval query",
         "Evaluate on larger, more complex document corpora",
     ], GREEN),
]

x = Inches(0.35)
for title, bullets, accent_col in next_steps:
    card(s, x, Inches(1.35), Inches(4.1), Inches(5.7))
    add_text(s, title, x + Inches(0.15), Inches(1.45),
             Inches(3.8), Inches(0.55),
             font_size=Pt(15), bold=True, color=accent_col)
    add_bullet_box(s, bullets, x + Inches(0.15), Inches(2.05),
                   Inches(3.8), Inches(4.5), font_size=Pt(13))
    x += Inches(4.31)

add_text(s, "Code available at: UEF-Research-Project  ·  GitHub · All experiments reproducible locally",
         Inches(0.35), Inches(7.1), Inches(12.63), Inches(0.35),
         font_size=Pt(13), color=TEAL, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out_path = "RAG_Research_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({prs.slides.__len__()} slides)")
